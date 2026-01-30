import base64
import re
from typing import Tuple, Optional
from io import BytesIO
from pypdf import PdfReader

from app.utils.ai_config import get_ai_client


def reflow_pdf_text(text: str) -> str:
    """
    Reflow PDF-extracted text to join lines that are part of the same paragraph.

    PDFs often have hard line breaks at visual line boundaries, causing text to appear
    as 1-3 words per line with blank lines between. This function joins such lines while
    preserving:
    - Real paragraph breaks (multiple blank lines or semantic breaks)
    - List items (lines starting with bullets, numbers, or letters)
    - Headings (short lines that look like titles)
    """
    if not text:
        return text

    # First pass: detect if this is a "word-per-line" PDF by checking the pattern
    lines = text.split('\n')
    non_blank_lines = [l.strip() for l in lines if l.strip()]

    # Calculate average words per non-blank line
    if non_blank_lines:
        total_words = sum(len(l.split()) for l in non_blank_lines)
        avg_words = total_words / len(non_blank_lines)
    else:
        return text

    # If average is very low (< 3 words per line), this is likely word-per-line extraction
    # In this case, we need to be much more aggressive about joining
    is_word_per_line = avg_words < 3

    if is_word_per_line:
        # For word-per-line PDFs: join everything, then split on semantic boundaries
        return _reflow_word_per_line(text)
    else:
        # For normal PDFs: use standard paragraph-based reflow
        return _reflow_normal(text)


def _reflow_word_per_line(text: str) -> str:
    """Handle PDFs where text is extracted as one word per line."""
    # Join all text first, treating blank lines as spaces
    # Then re-split on semantic boundaries

    lines = text.split('\n')
    tokens = []
    blank_count = 0

    for line in lines:
        stripped = line.strip()
        if not stripped:
            blank_count += 1
            continue

        # Multiple blanks (3+) might indicate a real paragraph break
        if blank_count >= 3:
            tokens.append('\n\n')
        blank_count = 0

        tokens.append(stripped)

    # Join with spaces
    joined = ' '.join(t if t != '\n\n' else t for t in tokens)
    # Clean up: paragraph markers become actual breaks
    joined = joined.replace(' \n\n ', '\n\n').replace('\n\n ', '\n\n').replace(' \n\n', '\n\n')
    # Clean up multiple spaces
    joined = re.sub(r' +', ' ', joined)

    # Now process line by line to restore structure
    paragraphs = joined.split('\n\n')
    result_parts = []

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        # Check for list items and split them out
        # Look for bullet patterns in the middle of text
        # Split on bullets: •, -, numbered items
        parts = re.split(r'(?=\s*[•●○▪▸]\s+)|(?=\s+\d+[\.\)]\s+)', para)

        for part in parts:
            part = part.strip()
            if not part:
                continue

            # Check if this starts with a numbered section like "1)" or "2)"
            if re.match(r'^\d+\)\s+', part):
                result_parts.append('')  # Add paragraph break before numbered sections
                result_parts.append(part)
            elif part.startswith(('•', '●', '○', '▪', '▸')):
                result_parts.append(part)
            else:
                result_parts.append(part)

    # Join and clean up
    result = '\n\n'.join(result_parts)
    result = re.sub(r'\n{3,}', '\n\n', result)

    return result.strip()


def _reflow_normal(text: str) -> str:
    """Handle normal PDFs with reasonable line lengths."""
    lines = text.split('\n')
    result_lines = []
    current_paragraph = []

    def is_list_item(line: str) -> bool:
        """Check if line starts with a list marker."""
        stripped = line.strip()
        if stripped.startswith(('•', '-', '*', '–', '—', '○', '●', '▪', '▸')):
            return True
        if re.match(r'^[\(\[]?\d+[\.\)\]]', stripped):
            return True
        if re.match(r'^[\(\[]?[a-zA-Z][\.\)\]]', stripped):
            return True
        return False

    def is_likely_heading(line: str, prev_blank: bool) -> bool:
        """Check if line looks like a heading."""
        stripped = line.strip()
        if not stripped:
            return False
        if prev_blank and len(stripped) < 80 and not stripped.endswith((',', ';')):
            if stripped[0].isupper() and not stripped.endswith(('.', ':', '?', '!')):
                return True
        if stripped.isupper() and len(stripped) < 60:
            return True
        return False

    def flush_paragraph():
        if current_paragraph:
            joined = ' '.join(current_paragraph)
            joined = re.sub(r' +', ' ', joined)
            result_lines.append(joined)
            current_paragraph.clear()

    prev_was_blank = True

    for line in lines:
        stripped = line.strip()

        if not stripped:
            flush_paragraph()
            if result_lines and result_lines[-1] != '':
                result_lines.append('')
            prev_was_blank = True
            continue

        if is_list_item(stripped):
            flush_paragraph()
            result_lines.append(stripped)
            prev_was_blank = False
            continue

        if is_likely_heading(stripped, prev_was_blank):
            flush_paragraph()
            result_lines.append(stripped)
            prev_was_blank = False
            continue

        if current_paragraph:
            last_para_text = current_paragraph[-1]
            ends_sentence = last_para_text.rstrip().endswith(('.', '!', '?', ':'))
            starts_capital = stripped[0].isupper() if stripped else False
            if ends_sentence and starts_capital and len(last_para_text) < 50:
                flush_paragraph()

        current_paragraph.append(stripped)
        prev_was_blank = False

    flush_paragraph()

    result = '\n'.join(result_lines)
    result = re.sub(r'\n{3,}', '\n\n', result)

    return result.strip()


class ArtifactConversionService:
    """Service for converting various file types to markdown using AI."""

    def __init__(self):
        self.client = get_ai_client()

    async def convert_to_markdown(
        self,
        file_content: bytes,
        filename: str,
        mime_type: str,
    ) -> Tuple[str, bool]:
        """
        Convert file content to markdown.

        Returns:
            Tuple of (markdown_content, success)
        """
        try:
            # Route to appropriate converter based on mime type
            if mime_type.startswith("image/"):
                return await self._convert_image(file_content, filename, mime_type)
            elif mime_type == "application/pdf":
                return await self._convert_pdf(file_content, filename)
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword",
            ):
                return await self._convert_docx(file_content, filename)
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            ):
                return await self._convert_xlsx(file_content, filename)
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint",
            ):
                return await self._convert_pptx(file_content, filename)
            elif mime_type.startswith("text/") or mime_type in (
                "application/json",
                "application/xml",
                "application/javascript",
            ):
                return await self._convert_text(file_content, filename, mime_type)
            else:
                return await self._convert_unknown(file_content, filename, mime_type)
        except Exception as e:
            return f"# Conversion Error\n\nFailed to convert `{filename}`: {str(e)}", False

    async def _convert_image(
        self, file_content: bytes, filename: str, mime_type: str
    ) -> Tuple[str, bool]:
        """Convert image to markdown with AI-generated description."""
        base64_content = base64.b64encode(file_content).decode("utf-8")

        prompt = """Analyze this image and create a detailed markdown description.

If this is a document or specification:
- Extract all text content
- Preserve any structure (headings, lists, tables)
- Describe any diagrams, charts, or visual elements

If this is a diagram or wireframe:
- Describe the overall purpose
- List all components and their relationships
- Note any labels, text, or annotations

If this is a screenshot:
- Describe what the interface shows
- Note any important UI elements
- Capture any visible text

Format your response as clean markdown with appropriate headings."""

        try:
            text = await self.client.complete(
                use_case="file_to_markdown",
                system_prompt="You are an expert at analyzing images and converting them to detailed markdown descriptions.",
                user_content=prompt,
                images=[{"media_type": mime_type, "data": base64_content}],
            )

            markdown = f"# {filename}\n\n{text}"
            return markdown, True

        except Exception as e:
            error_msg = str(e).lower()
            if "rate limit" in error_msg or "rate_limit" in error_msg:
                return (
                    f"# {filename}\n\n*AI service rate limit exceeded. Please wait a moment and try again.*",
                    False,
                )
            elif "authentication" in error_msg or "api key" in error_msg:
                return (
                    f"# {filename}\n\n*AI service authentication failed. Please check the API key configuration.*",
                    False,
                )
            elif "image" in error_msg or "media" in error_msg:
                return (
                    f"# {filename}\n\n*Invalid image format or size. Please ensure images are JPEG, PNG, GIF, or WebP and under 5MB each.*",
                    False,
                )
            return (
                f"# {filename}\n\n*AI service error: {str(e)}*",
                False,
            )

    async def _convert_pdf(
        self, file_content: bytes, filename: str
    ) -> Tuple[str, bool]:
        """Convert PDF to markdown, using Vision for image-heavy pages."""
        try:
            reader = PdfReader(BytesIO(file_content))
            text_parts = []
            has_content = False

            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    has_content = True
                    # Reflow the text to join lines that are part of the same paragraph
                    reflowed_text = reflow_pdf_text(page_text)
                    text_parts.append(f"## Page {i + 1}\n\n{reflowed_text}")

            if not has_content:
                # PDF might be image-based, use Vision
                return await self._convert_pdf_with_vision(file_content, filename)

            markdown = f"# {filename}\n\n" + "\n\n".join(text_parts)
            return markdown, True

        except Exception as e:
            return f"# {filename}\n\n*PDF extraction failed: {str(e)}*", False

    async def _convert_pdf_with_vision(
        self, file_content: bytes, filename: str
    ) -> Tuple[str, bool]:
        """Convert image-based PDF using Claude Vision."""
        # For now, return a placeholder - full PDF-to-image conversion
        # would require additional dependencies like pdf2image/poppler
        return (
            f"# {filename}\n\n*This PDF appears to be image-based. "
            "Text extraction was not possible. Consider uploading individual page images.*",
            False,
        )

    async def _convert_docx(
        self, file_content: bytes, filename: str
    ) -> Tuple[str, bool]:
        """Convert Word document to markdown."""
        try:
            from docx import Document

            doc = Document(BytesIO(file_content))
            parts = [f"# {filename}"]

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Handle headings
                if para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading", "").strip()
                    try:
                        level_num = int(level) + 1  # Convert to markdown heading
                    except ValueError:
                        level_num = 2
                    parts.append(f"{'#' * level_num} {text}")
                else:
                    parts.append(text)

            # Handle tables
            for table in doc.tables:
                table_md = self._table_to_markdown(table)
                if table_md:
                    parts.append(table_md)

            markdown = "\n\n".join(parts)
            return markdown, True

        except ImportError:
            return (
                f"# {filename}\n\n*python-docx not installed. Cannot convert Word documents.*",
                False,
            )
        except Exception as e:
            return f"# {filename}\n\n*Word document conversion failed: {str(e)}*", False

    def _table_to_markdown(self, table) -> Optional[str]:
        """Convert a docx table to markdown format."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")

        if not rows:
            return None

        # Add header separator after first row
        if len(rows) > 0:
            col_count = len(table.rows[0].cells)
            separator = "| " + " | ".join(["---"] * col_count) + " |"
            rows.insert(1, separator)

        return "\n".join(rows)

    async def _convert_xlsx(
        self, file_content: bytes, filename: str
    ) -> Tuple[str, bool]:
        """Convert Excel spreadsheet to markdown tables."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(BytesIO(file_content), data_only=True)
            parts = [f"# {filename}"]

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                parts.append(f"## {sheet_name}")

                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Skip empty rows
                    if all(cell is None for cell in row):
                        continue
                    cells = [
                        str(cell if cell is not None else "").replace("|", "\\|")
                        for cell in row
                    ]
                    rows.append("| " + " | ".join(cells) + " |")

                if rows:
                    # Add header separator
                    col_count = len(rows[0].split("|")) - 2
                    separator = "| " + " | ".join(["---"] * col_count) + " |"
                    rows.insert(1, separator)
                    parts.append("\n".join(rows))
                else:
                    parts.append("*Empty sheet*")

            markdown = "\n\n".join(parts)
            return markdown, True

        except ImportError:
            return (
                f"# {filename}\n\n*openpyxl not installed. Cannot convert Excel files.*",
                False,
            )
        except Exception as e:
            return (
                f"# {filename}\n\n*Excel conversion failed: {str(e)}*",
                False,
            )

    async def _convert_pptx(
        self, file_content: bytes, filename: str
    ) -> Tuple[str, bool]:
        """Convert PowerPoint presentation to markdown."""
        try:
            from pptx import Presentation

            prs = Presentation(BytesIO(file_content))
            parts = [f"# {filename}"]

            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"## Slide {i}")
                slide_texts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                if slide_texts:
                    parts.append("\n\n".join(slide_texts))
                else:
                    parts.append("*No text content*")

            markdown = "\n\n".join(parts)
            return markdown, True

        except ImportError:
            return (
                f"# {filename}\n\n*python-pptx not installed. Cannot convert PowerPoint files.*",
                False,
            )
        except Exception as e:
            return (
                f"# {filename}\n\n*PowerPoint conversion failed: {str(e)}*",
                False,
            )

    async def _convert_text(
        self, file_content: bytes, filename: str, mime_type: str
    ) -> Tuple[str, bool]:
        """Convert text/code files to markdown with code blocks."""
        try:
            text = file_content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = file_content.decode("latin-1")
            except Exception:
                return (
                    f"# {filename}\n\n*Could not decode text content*",
                    False,
                )

        # Determine language for code block
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        lang_map = {
            "py": "python",
            "js": "javascript",
            "ts": "typescript",
            "tsx": "typescript",
            "jsx": "javascript",
            "json": "json",
            "xml": "xml",
            "html": "html",
            "css": "css",
            "sql": "sql",
            "sh": "bash",
            "yaml": "yaml",
            "yml": "yaml",
            "md": "markdown",
        }

        if ext in lang_map:
            lang = lang_map[ext]
            markdown = f"# {filename}\n\n```{lang}\n{text}\n```"
        elif mime_type == "text/plain" or ext == "txt":
            markdown = f"# {filename}\n\n{text}"
        else:
            markdown = f"# {filename}\n\n```\n{text}\n```"

        return markdown, True

    async def _convert_unknown(
        self, file_content: bytes, filename: str, mime_type: str
    ) -> Tuple[str, bool]:
        """Handle unknown file types."""
        size_kb = len(file_content) / 1024
        return (
            f"# {filename}\n\n"
            f"**File Type:** {mime_type}\n"
            f"**Size:** {size_kb:.1f} KB\n\n"
            "*This file type cannot be converted to markdown. "
            "The original file is available for download.*",
            False,
        )
